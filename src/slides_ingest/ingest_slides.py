import logging
import os
import re
from datetime import datetime
from typing import List, Optional, Tuple, Union

import pandas as pd
from pptx import Presentation
from tqdm import tqdm

from transformers import pipeline

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

log = logging.getLogger()

class SlidesIngest:
    """
    Class for PowerPoint slides content ingesting and summarisation
    """
    def __init__(
        self, 
        pp_filename: str
        ):
        """
        :param pp_filename: PowerPoint file location
        """
        self.pp_filename = pp_filename
        
        self.base_path = '/Users/ansonliu/Downloads'
        self.filepath = os.path.join(self.base_path, self.pp_filename)
        self.today = datetime.today().strftime('%d-%m-%Y')
        
        self.long_sum = None
        
        self.prs = None
        self.all_slides_text = None
        self.slide_content = None
        self.slide_summary = None

        self.references_df = None

    def init_summarisers(self) -> None:
        """
        Function to initialise all summariser tools
        """
        log.info('Initialisiing hugging face summary tools')
        
        # self.long_sum = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")
        self.long_sum = pipeline('summarization', model='facebook/bart-large-cnn')
    
    @staticmethod
    def calc_min_max_tokens(
        input_length: int
    ) -> Tuple[int, int]:
        """
        Function to calculate the minimum and maximum tokens to use/cap at
        :param input_length: numbers of words in string
        :return: min and max token values
        """
        min_length = max(1, int(input_length * 0.1))
        max_length = max(5, min(input_length // 1.25, 200))

        return min_length, max_length

    def load_file(self) -> None:
        """
        Function to load powerpoint and extract content.
        """
        log.info(f'Loading PowerPoint file: {self.filepath}')
        
        self.slide_content = {}
        self.prs = Presentation(self.filepath)

        log.info(f'PowerPoint total slide count: {len(self.prs.slides)}')
        
        for slide_num, slide in enumerate(self.prs.slides, start=0):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_text.append(shape.text)
            
            self.slide_content[slide_num] = " ".join(slide_text) if slide_text else ""

        log.info(f'Successfully extracted content from {len(self.slide_content)} slides.')

    def get_slide_text(
        self,
        slide_mapped: bool,
        n_slides: Optional[Union[int, List[int]]] = None
    ) -> str:
        """
        Function to get slides from all or a specific set of slides 
        :param slide_mapped: Option to show which summary is from which slide
        :param n_slides: slide selection range, if None will get all
        :return all_slide_text: text from slide
        """
        log.info(f'Getting{" all " if n_slides is None else " "}slide text{f"" if n_slides is None else " in slides {n_slides}"}')
        
        all_slides = True if n_slides is None else False
        
        n_slides = n_slides if n_slides else [n for n in range(len(self.prs.slides))]
        n_slides = [n_slides] if isinstance(n_slides, int) else n_slides

        all_slide_text = f'. '.join([f'{f"Slide {n + 1}: " if slide_mapped else ""} {self.slide_content[n]}' for n in n_slides if self.slide_content[n]])

        return all_slide_text

    @staticmethod
    def split_text_chunks(
        text: str, 
        max_token_length: int = 200
        ) -> List[str]:
        """
        Splits text into chunks that do not exceed the model's maximum token limit.
        :param text: Input text to split
        :param max_token_length: Token limit per chunk
        :return: List of text chunks
        """
        words = text.split()
        return [" ".join(words[i:i + max_token_length]) for i in range(0, len(words), max_token_length)]

    def summarise_all(self) -> None:
        """
        Function to summarise all slides as one, splitting into chunks if input exceeds model's token limit.
        """
        log.info('Summarising all slides')

        self.all_slides_text = self.get_slide_text(slide_mapped=False)
        input_length = len(self.all_slides_text.split())

        # Check if text exceeds token limit
        if input_length > 200:
            log.info("Splitting text into smaller chunks due to token limit")
            text_chunks = self.split_text_chunks(self.all_slides_text)
            chunk_summaries = []
            
            n_chunk = 1

            for chunk in tqdm(text_chunks, leave=True):
                log.info(f'Summarising chunk {n_chunk} of total {len(text_chunks)}')
                n_chunk += 1
                
                min_length, max_length = self.calc_min_max_tokens(input_length=len(chunk.split()))
                summary = self.long_sum(
                    chunk,
                    max_length=max_length,
                    min_length=min_length,
                    do_sample=False
                )[0]['summary_text']
                chunk_summaries.append(summary)
            
            combined_text = " ".join(chunk_summaries)
            
            self.slide_summary = combined_text
        else:
            # Summarize directly if within token limit
            min_length, max_length = self.calc_min_max_tokens(input_length=input_length)
            pp_summary = self.long_sum(
                self.all_slides_text,
                max_length=max_length,
                min_length=min_length,
                do_sample=False
            )
            self.slide_summary['all'] = pp_summary[0]['summary_text']
    
    def get_references(self) -> None:
        """
        Function to get all references from slides
        """
        log.info('Extracting references')

        references_patterns = [
            r'([A-Za-z]+ \(\d{4}\))',                 # Name (yyyy)
            r'(\([A-Za-z]+, \d{4}\))',                # (Name, yyyy)
            r'([A-Za-z]+ and [A-Za-z]+ \(\d{4}\))',   # Name and Name (yyyy)
            r'(\([A-Za-z]+ and [A-Za-z]+, \d{4}\))',  # (Name and Name, yyyy)
            r'([A-Za-z]+ et al\., \(\d{4}\))',        # Name et al. (yyyy)
            r'(A-Za-z]+ et al\., yyyy)\)'             # (Name et al., yyyy)
        ]
        
        self.references_df = pd.DataFrame(columns=['References', 'Slide'])
        
        for pattern in references_patterns:
            references = re.findall(pattern, self.all_slides_text)
            
            if references:
                self.format_references(references) 
        
        self.references_df = self.references_df.sort_values(by=['References', 'Slide']).drop_duplicates().reset_index(drop=True)

    def format_references(self, references: List[str]) -> None:
        """
        Function to format references and locate which slide it was presented
        :param references: list of references found
        """
        for reference in references:
            for slide_n, slide_text in self.slide_content.items():
                if reference in slide_text:
                    self.references_df = pd.concat([self.references_df, pd.DataFrame({'Slide': [slide_n], 'References': [reference]})], ignore_index=True).reset_index(drop=True)
            
            if reference not in self.references_df['References'].unique():
                self.references_df = pd.concat([self.references_df, pd.DataFrame({'Slide': ['No slide number found'], 'References': [reference]})], ignore_index=True).reset_index(drop=True)
    
    def display_summary(self) -> None:
        """
        Function to display summarisation in terminal
        """
        sum_list = [point for point in self.slide_summary.split('. ')]
        
        print('-------------------- SUMMARISATION START --------------------')
        for n_point, point in enumerate(sum_list, start=1):
            print(f'{n_point}. {point}')
        print('------------------------ REFERENCES -------------------------')
        print(self.references_df.to_string())
        print('--------------------- SUMMARISATION END ---------------------')
        
    def save_summarisation(self) -> None:
        """
        Function to save summarisation to documents
        """
        raise NotImplementedError('Not implemented')
        # save_filename = f'{self.pp_filename.split('.')[0]}_summarised.txt'
        # save_loc = os.path.join(self.base_path, save_filename)
        # log.info(f'Saving summarisation: {save_loc}')
        
    def summarise_runner(self) -> None: 
        """
        Sub-runner function for all summarisation processes
        """
        self.summarise_all()
        self.get_references()
        
    def runner(
        self,
        save: bool,
        display: bool
        ) -> None:
        """
        Main runner function
        :param save: Option to save to csv
        :param display: Option to display summarisations
        """
        log.info(f'Hello, Danielle. Current date: {self.today}')
        log.info(f'Starting PowerPoint summarisation on file: {self.pp_filename}')
        
        self.init_summarisers()
        self.load_file()
        
        self.summarise_runner()
        
        if display:
            self.display_summary()
        
        if save: 
            self.save_summarisation()
        
        log.info('Summarisation Complete')
        